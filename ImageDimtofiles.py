import os
import io
import time
from google import genai
from google.genai import types
import google.genai.errors
from pydantic import BaseModel, Field
from PIL import Image
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv

# ==========================================
# 1. INITIALIZATION & DYNAMIC TOGGLE SYSTEM
# ==========================================
load_dotenv()

# Read the current active environment mode switch from the hidden vault file
current_env_mode = os.environ.get("CURRENT_ENV", "GEMINI").strip().upper()

# Select the appropriate API key based on the environment toggle
if current_env_mode == "GROQ":
    active_api_key = os.environ.get("GROQ_API_KEY")
    print("🔌 System Engine Booted using: GROQ Cloud Environment Key")
elif current_env_mode == "GEMINI":
    active_api_key = os.environ.get("GEMINI_API_KEY")
    print("🔌 System Engine Booted using: GOOGLE Gemini Environment Key")
else:
    print("⚠️ Unknown configuration value in CURRENT_ENV. Defaulting to Gemini.")
    active_api_key = os.environ.get("GEMINI_API_KEY")

# Initialize the primary Client with the chosen operational key
client = genai.Client(api_key=active_api_key)

# ==========================================
# 2. STRICT DATA STRUCTURING (SCHEMA)
# ==========================================
class ObjectMeasurement(BaseModel):
    object_name: str = Field(description="Name of object (e.g., Tree, Circle, Lens)")
    shape_type: str = Field(description="Geometry shape (e.g., Circular, Rectangular, Organic)")
    estimated_width: str = Field(description="Real-world width with units")
    estimated_height: str = Field(description="Real-world height with units")
    estimated_area: str = Field(description="Calculated real-world surface footprint area with squared units (e.g., sq meters, sq cm)")
    confidence_and_notes: str = Field(description="Brief estimation context and geometric math choice used")

class ImageAnalysisResult(BaseModel):
    items: list[ObjectMeasurement]

# ==========================================
# 3. HIGH-EFFICIENCY VISION ENGINE WITH RETRY
# ==========================================
def analyze_image_dimensions(image_path):
    """
    Compresses target image dynamically to reduce upload latency,
    and streams payload to the API with automated retry logic for 503 errors.
    """
    if not os.path.exists(image_path):
        print(f"❌ File not found: {image_path}")
        return None

    print(f"\n⚡ Optimizing & compressing image memory footprint...")
    try:
        with Image.open(image_path) as img:
            # Downscale massive images to a max layout boundary of 1024px to save transmission bandwidth
            img.thumbnail((1024, 1024))
            
            # Save into temporary system RAM (Byte Buffer) instead of making slow files on your hard drive
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='JPEG', quality=85)
            img_bytes = img_byte_arr.getvalue()
    except Exception as e:
        print(f"❌ Error optimizing image data locally: {e}")
        return None

    prompt = (
        "Analyze this image. Detect primary objects (shapes, trees, products, components). "
        "Estimate real-world width and height using background context clues or scale indicators. "
        "Calculate the real-world surface footprint area based on its geometric boundaries. "
        "Strictly follow the requested JSON data schema structure."
    )

    max_retries = 3
    for attempt in range(max_retries):
        try:
            print(f"🛰️ Streaming binary stream to Vision AI Engine (Attempt {attempt + 1}/{max_retries})...")
            response = client.models.generate_content(
                model='gemini-2.5-flash',
                contents=[
                    types.Part.from_bytes(data=img_bytes, mime_type='image/jpeg'),
                    prompt
                ],
                config=types.GenerateContentConfig(
                    response_mime_type="application/json",
                    response_schema=ImageAnalysisResult,
                    temperature=0.1
                ),
            )
            return response.parsed  # Break retry loop and return successfully parsed structural map
            
        except google.genai.errors.ServerError as e:
            if "503" in str(e) and attempt < max_retries - 1:
                print("⚠️ API gateway servers are heavily loaded. Pausing 3 seconds before auto-retry...")
                time.sleep(3)
            else:
                print(f"❌ Server Error returned from Cloud Provider: {e}")
                return None
        except Exception as e:
            print(f"❌ Unexpected processing error occurred: {e}")
            return None

# ==========================================
# 4. EXPORT ENGINE SUITE
# ==========================================
def export_to_excel(analysis_data, filename="Detected_Dimensions.xlsx"):
    wb = Workbook()
    ws = wb.active
    ws.title = "Spatial Summary"
    
    # 6-Column Structured Header Map
    ws.append(["Object Name", "Shape Type", "Width", "Height", "Estimated Area", "Analysis Notes"])
    
    for item in analysis_data.items:
        ws.append([
            item.object_name, 
            item.shape_type, 
            item.estimated_width, 
            item.estimated_height, 
            item.estimated_area, 
            item.confidence_and_notes
        ])
        
    wb.save(filename)
    print(f"📊 Excel Ledger Successfully Saved: '{filename}'")

def export_to_docx(analysis_data, filename="Detected_Dimensions.docx"):
    doc = Document()
    doc.add_heading('AI Spatial Measurement & Area Report', level=0)
    
    # 5-Column Document Table Layout Matrix
    table = doc.add_table(rows=1, cols=5)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Object'
    hdr_cells[1].text = 'Geometry'
    hdr_cells[2].text = 'Width'
    hdr_cells[3].text = 'Height'
    hdr_cells[4].text = 'Estimated Area'
    
    for item in analysis_data.items:
        row_cells = table.add_row().cells
        row_cells[0].text = str(item.object_name)
        row_cells[1].text = str(item.shape_type)
        row_cells[2].text = str(item.estimated_width)
        row_cells[3].text = str(item.estimated_height)
        row_cells[4].text = str(item.estimated_area)
        
        p = doc.add_paragraph()
        p.add_run(f"\n🔍 {item.object_name} Analytical Context: ").bold = True
        p.add_run(item.confidence_and_notes)
        
    doc.save(filename)
    print(f"📝 Word Document Summary Successfully Saved: '{filename}'")

def export_to_pptx(analysis_data, filename="Detected_Dimensions.pptx"):
    prs = Presentation()
    
    # --- SLIDE 1: TITLE SLIDE ---
    title_slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI Spatial Measurement Report"
    subtitle.text = f"Automated Dimension & Area Extraction\nGenerated on: {time.strftime('%Y-%m-%d')}"
    
    # --- SLIDE 2+: DATA SLIDES (One slide per detected object) ---
    for item in analysis_data.items:
        bullet_slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Set Slide Header
        slide.shapes.title.text = f"Detected Object: {item.object_name}"
        
        # Populate text bullet points frame
        tf = slide.placeholders[1].text_frame
        tf.text = f"• Geometry Shape: {item.shape_type}"
        
        p2 = tf.add_paragraph()
        p2.text = f"• Estimated Width: {item.estimated_width}"
        
        p3 = tf.add_paragraph()
        p3.text = f"• Estimated Height: {item.estimated_height}"
        
        p4 = tf.add_paragraph()
        p4.text = f"• Calculated Area: {item.estimated_area}"
        p4.font.bold = True  # High contrast focus
        
        p5 = tf.add_paragraph()
        p5.text = f"• Analysis Notes: {item.confidence_and_notes}"
        
    prs.save(filename)
    print(f"📉 PowerPoint Deck Successfully Saved: '{filename}'")

# ==========================================
# 5. AUTOMATED INTERACTIVE INTERFACE
# ==========================================
if __name__ == "__main__":
    VALID_EXTENSIONS = ('.jpg', '.jpeg', '.png', '.webp')
    
    # Automatically scan active runtime workspace folder environment for asset targets
    current_directory = os.getcwd()
    all_files = os.listdir(current_directory)
    image_files = [f for f in all_files if f.lower().endswith(VALID_EXTENSIONS)]
    
    print("\n===========================================")
    print("🤖 VISIONAI: SPATIAL MEASUREMENT RUNTIME")
    print("===========================================")
    print(f"Active Engine Profile Mode: {current_env_mode}")
    print("===========================================")
    
    if not image_files:
        print("❌ No valid local images identified inside your project root folder!")
        print(f"Please drop evaluation assets ({', '.join(VALID_EXTENSIONS)}) directly into: {current_directory}")
    else:
        print("Detected project image assets available for analysis:")
        for index, img_name in enumerate(image_files, start=1):
            print(f"  [{index}] {img_name}")
            
        try:
            img_choice = int(input("\nSelect target asset entry number to evaluate: ").strip())
            if img_choice < 1 or img_choice > len(image_files):
                print("❌ Input bounds error. Terminating execution context.")
                exit()
            
            test_image_path = image_files[img_choice - 1]
            print(f"🎯 Operational Target Locked: '{test_image_path}'")
            
        except ValueError:
            print("❌ Numeric parser error. Integer entry required.")
            exit()

        print("\n--- 🎯 TARGET CHOOSE OUTPUT FORMAT ---")
        user_choice = input("Type 'E' for Excel, 'W' for Word, or 'P' for PowerPoint Presentation: ").strip().upper()
        
        if user_choice not in ['E', 'W', 'P']:
            print("❌ Output selector error. Execution terminated.")
        else:
            # Fire processing pipeline
            results = analyze_image_dimensions(test_image_path)
            
            if results and results.items:
                print(f"\n✅ Successfully processed spatial calculations for {len(results.items)} objects!")
                
                if user_choice == 'E':
                    export_to_excel(results, "Detected_Dimensions.xlsx")
                elif user_choice == 'W':
                    export_to_docx(results, "Detected_Dimensions.docx")
                elif user_choice == 'P':
                    export_to_pptx(results, "Detected_Dimensions.pptx")
                    
                print("🏁 Data pipeline successfully finalized. Happy Reviewing!")
            else:
                print("❌ System reached data endpoint without valid metric extractions.")